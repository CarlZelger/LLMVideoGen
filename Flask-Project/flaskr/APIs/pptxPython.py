import os
import subprocess
from typing import Dict, List
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def getPageByIndex(i: int,prs):
    title_slide_layout = prs.slide_layouts[i]
    slide = prs.slides.add_slide(title_slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    return slide

def setTitle(slide, title: str, left=Inches(0.5), top=Inches(0.5), width=Inches(8), height=Inches(1)):
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_shape.text_frame
    text_frame.text = title
    # Optionally set font size and other properties
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(32)  # Adjust the font size as needed
            run.font.name = 'Arial'
            run.font.color.rgb = RGBColor(255, 255, 255)  # Set font color to white

def setSubtitle(slide, subtitle: str, left=Inches(0.5), top=Inches(1.5), width=Inches(8), height=Inches(1)):
    subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = subtitle_shape.text_frame
    text_frame.text = subtitle
    # Optionally set font size and other properties
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)  # Adjust the font size as needed
            run.font.name = 'Arial'
            run.font.color.rgb = RGBColor(255, 255, 255)

def setNote(slide, note: str):
    n = slide.notes_slide
    text_frame = n.notes_text_frame
    text_frame.text = note
    
def addImage(slide, path):
    left = Inches(0)  # Position from the left of the slide
    top = Inches(2.27)  # Position from the top of the slide
    height = Inches(5.23)  # Height of the image
    
    top = Inches(1.2)  # Position from the top of the slide
    height = Inches(6.3)  # Height of the image
    
    slide.shapes.add_picture(path, left, top,height=height)
    # if os.path.exists(path):
    #     os.remove(path)
    
    
def addEndSlide(prs):
    slide = getPageByIndex(0,prs)
    setTitle(slide, "Add improvements!")
    node = 'I hope the video was helpful. we generated some recommendations for improvements, if you are interested: simply press the button below that matches your expectations for a remade video. Thanks for watching and I hope to see you again soon on our site'
    setNote(slide, node)
    
def generatePresentation(titles: List[str], content: Dict[str,str], images: Dict[str,str],question:str=""):
    prs = Presentation()
    fname = 'test1.pptx'
    for t in titles:
        slide = getPageByIndex(0,prs)
        setTitle(slide, t)
        setNote(slide, content[t])
        addImage(slide, images[t])
    if question is not "":
        print("adding Question")
        addQuestionSlide(question,prs)
        fname = 'test2.pptx'
    addEndSlide(prs)
        
    prs.save(fname)
    return fname

def addQuestionSlide(q: str,prs):
    slide = getPageByIndex(0,prs)
    setTitle(slide, "Ecercises")
    note = f"Here is a quick question for you to strenghten youre knowledge: {q}"
    setSubtitle(slide,q)
    setNote(slide,note)







